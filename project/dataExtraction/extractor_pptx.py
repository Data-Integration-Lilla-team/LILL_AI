from pptx import Presentation

class PPTX_Text_Extractor:

    def __init__(self):
        pass
    

    def sort_dic(self,dic):
        sorted_dict = {}
        for key in sorted(dic.keys()):
            sorted_dict[key] = dic[key]
        return sorted_dict


    def extract_data(self, list_of_path):
        print("Estrazione dati pptx")
        dict_path_page_text = {}
        for file in list_of_path:
            presentation = Presentation(file)
            for idx, slide in enumerate(presentation.slides):
                
                text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text
                key = file + "#" + str(idx)
                dict_path_page_text[key] = text
        return dict_path_page_text